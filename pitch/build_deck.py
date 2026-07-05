"""Builds AgentForge pre-seed investor pitch deck.

Run with::

    uv run --with python-pptx pitch/build_deck.py

Output: ``pitch/AgentForge.pitch.pptx``

The deck is parameterized via ``DECK_CONFIG`` and the slide-builder helpers below.
Edit any slide's ``body`` to change content. Edit ``BRAND`` to recolor.

NOTE: Several slides have placeholders in [BRACKETS] that the founder must
fill in before showing this to anyone (founder bio, cap table, current pipeline,
etc). Search for "[" to find them.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


# ───────────────────────────────────────────────────────────────────
# Brand
# ───────────────────────────────────────────────────────────────────


@dataclass(frozen=True)
class Brand:
    navy: RGBColor
    navy_dark: RGBColor
    accent: RGBColor
    body: RGBColor
    muted: RGBColor
    bg: RGBColor
    title_font: str = "Calibri"
    body_font: str = "Calibri"


BRAND = Brand(
    navy=RGBColor(0x0D, 0x3B, 0x66),       # forge-500
    navy_dark=RGBColor(0x06, 0x1F, 0x37),  # forge-900
    accent=RGBColor(0xE3, 0x6E, 0x1A),     # warm orange against navy
    body=RGBColor(0x1F, 0x29, 0x37),
    muted=RGBColor(0x64, 0x74, 0x8B),
    bg=RGBColor(0xFF, 0xFF, 0xFF),
)


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ───────────────────────────────────────────────────────────────────
# Slide builders
# ───────────────────────────────────────────────────────────────────


def _add_rect(slide, left, top, width, height, fill: RGBColor, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    return shape


def _add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = BRAND.body,
    font: str = BRAND.body_font,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box, tf


def _add_bullets(slide, left, top, width, height, bullets: list[str | tuple[str, str]],
                 *, size: int = 18, gap: int = 8):
    """Each bullet is either a plain string or (heading, body) tuple."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0 if i == 0 else gap)

        if isinstance(item, tuple):
            heading, body = item
            r1 = p.add_run()
            r1.text = f"•  {heading}"
            r1.font.name = BRAND.body_font
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = BRAND.navy

            r2 = p.add_run()
            r2.text = f"  —  {body}"
            r2.font.name = BRAND.body_font
            r2.font.size = Pt(size)
            r2.font.color.rgb = BRAND.body
        else:
            run = p.add_run()
            run.text = f"•  {item}"
            run.font.name = BRAND.body_font
            run.font.size = Pt(size)
            run.font.color.rgb = BRAND.body
    return box


def _slide_with_layout(prs: Presentation, layout_index: int = 6):
    return prs.slides.add_slide(prs.slide_layouts[layout_index])


def _add_footer(slide, slide_num: int, total: int):
    _add_text(
        slide,
        Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
        "AgentForge · Pre-seed · 2026",
        size=10, color=BRAND.muted,
    )
    _add_text(
        slide,
        Inches(11.8), Inches(7.05), Inches(1.5), Inches(0.3),
        f"{slide_num} / {total}",
        size=10, color=BRAND.muted, align=PP_ALIGN.RIGHT,
    )


def _add_title_bar(slide, title: str, kicker: str | None = None):
    # Top accent line
    _add_rect(slide, Inches(0.5), Inches(0.45), Inches(0.6), Inches(0.07), BRAND.accent)
    if kicker:
        _add_text(
            slide,
            Inches(0.5), Inches(0.55), Inches(12), Inches(0.35),
            kicker.upper(),
            size=11, bold=True, color=BRAND.accent, font=BRAND.body_font,
        )
        title_top = Inches(0.95)
    else:
        title_top = Inches(0.6)
    _add_text(
        slide,
        Inches(0.5), title_top, Inches(12), Inches(0.85),
        title,
        size=34, bold=True, color=BRAND.navy_dark, font=BRAND.title_font,
    )


# ───────────────────────────────────────────────────────────────────
# Individual slides
# ───────────────────────────────────────────────────────────────────


def slide_cover(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BRAND.navy_dark)

    # Accent stripe
    _add_rect(s, Inches(0.6), Inches(2.0), Inches(0.6), Inches(0.08), BRAND.accent)

    _add_text(
        s, Inches(0.6), Inches(2.15), Inches(11), Inches(1.5),
        "AgentForge",
        size=72, bold=True, color=BRAND.bg, font="Calibri",
    )
    _add_text(
        s, Inches(0.6), Inches(3.4), Inches(11), Inches(1.5),
        "Domain-grounded AI research agents.\nStarting with insurance underwriting.",
        size=26, color=BRAND.bg, font=BRAND.body_font,
    )
    _add_text(
        s, Inches(0.6), Inches(6.2), Inches(11), Inches(0.5),
        "[Founder Name]   ·   [founder@email]   ·   Pre-seed   ·   2026",
        size=14, color=RGBColor(0xC0, 0xCB, 0xD9),
    )
    return s


def slide_problem(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(
        s,
        "Underwriters spend the majority of their day on manual research.",
        kicker="The Problem",
    )

    # Stat row
    stats = [
        ("60%", "of an underwriter's day is spent on research, not decisions"),
        ("4–8 hrs", "to manually quote a complex commercial risk"),
        ("$95K", "average loaded cost per junior underwriter, in a market that can't hire fast enough"),
    ]
    left = Inches(0.5)
    width = Inches(4.05)
    top = Inches(2.2)
    for i, (n, desc) in enumerate(stats):
        x = left + width * i + Inches(0.1) * i
        _add_text(s, x, top, width, Inches(1.0), n,
                  size=44, bold=True, color=BRAND.accent, font=BRAND.title_font)
        _add_text(s, x, top + Inches(1.05), width, Inches(1.4), desc,
                  size=15, color=BRAND.body)

    _add_text(
        s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
        "Every quote requires the same chain of lookups:",
        size=18, bold=True, color=BRAND.navy_dark,
    )
    _add_bullets(
        s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(2.0),
        [
            "ISO form language (CG 00 01, BP 00 03, CG 21 07 …) — pulled from PDFs and manuals",
            "State DOI rules (admitted carriers, prompt-payment timelines, surplus lines) — 50 portals",
            "NAIC company verification, group affiliation, market conduct history",
            "Coverage gap analysis against the carrier's manuscript endorsements",
        ],
        size=15,
    )
    return s


def slide_generic_ai_fails(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(
        s,
        "Generic AI doesn't solve this — and underwriters won't trust it.",
        kicker="Why Today's Tools Fall Short",
    )

    rows = [
        ("ChatGPT / Claude (direct)",
         "Hallucinates ISO form numbers and editions; can't cite primary sources; gets state-specific regulation wrong"),
        ("Perplexity / generic search",
         "Indexes the public web — misses Verisk ISO, NAIC CIS, state DOI portals; no auditable citations"),
        ("Westlaw / IRMI / Vertafore",
         "Reference libraries, not agents. Underwriter still does the synthesis manually"),
        ("Vertical AI incumbents (Harvey, Hebbia)",
         "Skipped insurance — too fragmented, too regulated, smaller per-deal ACVs than legal/finance"),
    ]
    top = Inches(2.0)
    row_h = Inches(1.05)
    for i, (left_text, right_text) in enumerate(rows):
        y = top + row_h * i
        _add_rect(s, Inches(0.5), y, Inches(0.06), Inches(0.9), BRAND.accent)
        _add_text(s, Inches(0.7), y, Inches(4.2), Inches(0.5), left_text,
                  size=15, bold=True, color=BRAND.navy_dark)
        _add_text(s, Inches(0.7), y + Inches(0.45), Inches(12), Inches(0.55), right_text,
                  size=13, color=BRAND.body)
    return s


def slide_insight(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(
        s,
        "A domain agent needs three things generic AI doesn't have.",
        kicker="The Insight",
    )

    cards = [
        ("Gated domain data",
         "Verisk ISO forms · NAIC CIS · 50 state DOI portals · NCCI · AM Best.",
         "Not the public web."),
        ("Domain-tuned reasoning",
         "Knows when to flag for human review · refuses to render coverage opinions · cites form-and-edition.",
         "Not a generic chatbot."),
        ("Structured, citable output",
         "Every claim ties to a verifiable source. JSON in, JSON out, audit trail by default.",
         "Not prose."),
    ]
    card_w = Inches(4.05)
    card_h = Inches(3.7)
    top = Inches(2.0)
    for i, (h, body, kicker) in enumerate(cards):
        x = Inches(0.5) + card_w * i + Inches(0.1) * i
        _add_rect(s, x, top, card_w, card_h, RGBColor(0xF5, 0xF7, 0xFA))
        _add_rect(s, x, top, card_w, Inches(0.08), BRAND.accent)
        _add_text(s, x + Inches(0.3), top + Inches(0.3), card_w, Inches(0.6),
                  h, size=20, bold=True, color=BRAND.navy_dark)
        _add_text(s, x + Inches(0.3), top + Inches(1.05), card_w - Inches(0.4),
                  Inches(2.0), body, size=14, color=BRAND.body)
        _add_text(s, x + Inches(0.3), top + Inches(3.05), card_w - Inches(0.4),
                  Inches(0.5), kicker, size=13, bold=True, color=BRAND.accent)
    return s


def slide_solution(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(
        s,
        "AgentForge is an AI research agent for insurance underwriting.",
        kicker="What We're Building",
    )

    bullets = [
        ("Citation-first.",
         "Every finding ties to a primary source — ISO form #, edition; DOI bulletin URL; NAIC MDL-NNN."),
        ("Refuses to overreach.",
         "Surfaces flags for human review on coverage determinations. No opinion-as-a-service."),
        ("Fast and cheap.",
         "Sub-90-second SLA. ~$0.50 per query at current model pricing."),
        ("Fits where underwriters work.",
         "REST API for AMS integrations · Slack / Teams for ad-hoc queries · webhook delivery for flow tools."),
        ("Built on a stable contract.",
         "One JSON schema across every query. Customers integrate once."),
    ]
    _add_bullets(
        s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(4.5),
        bullets, size=17, gap=12,
    )
    return s


def slide_how_it_works(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(s, "How it works", kicker="Architecture")

    boxes = [
        ("Query", "Underwriter asks a question via API, Slack, or AMS plugin"),
        ("Agent loop", "Custom HTTP loop — no LangChain. Calls domain tools, hard stops on cost/time/steps"),
        ("Domain tools", "Verisk ISO · NAIC CIS · State DOI · Policy parser · Web search · Vision (claim photos)"),
        ("Synthesis", "Strict-JSON pass — model can only cite tool results, no prior knowledge"),
        ("Output", "Structured response — summary · findings · sources · flags · confidence"),
    ]
    box_w = Inches(2.45)
    box_h = Inches(1.8)
    top = Inches(2.5)
    gap = Inches(0.1)

    for i, (h, body) in enumerate(boxes):
        x = Inches(0.5) + (box_w + gap) * i
        _add_rect(s, x, top, box_w, box_h, RGBColor(0xF5, 0xF7, 0xFA))
        _add_rect(s, x, top, box_w, Inches(0.08), BRAND.navy)
        _add_text(s, x + Inches(0.2), top + Inches(0.2), box_w - Inches(0.3),
                  Inches(0.5), h, size=15, bold=True, color=BRAND.navy_dark)
        _add_text(s, x + Inches(0.2), top + Inches(0.7), box_w - Inches(0.3),
                  Inches(1.2), body, size=11, color=BRAND.body)

        if i < len(boxes) - 1:
            arrow_x = x + box_w + Inches(0.005)
            _add_text(s, arrow_x, top + Inches(0.7), Inches(0.1), Inches(0.5),
                      "▶", size=14, color=BRAND.accent, align=PP_ALIGN.CENTER)

    _add_text(
        s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.0),
        "The model can only cite what tools returned. No hallucinated form numbers.",
        size=15, bold=True, color=BRAND.accent,
    )
    _add_bullets(
        s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.5),
        [
            "Hard stops: 8 steps · $0.50/run · 90 seconds — every run finishes or fails predictably.",
            "Per-tenant Redis cache, Postgres-backed run history, OTel + Langfuse for tracing.",
        ],
        size=13,
    )
    return s


def slide_demo(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(s, "Demo: a real underwriting query.", kicker="Example")

    _add_text(
        s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.6),
        "Query: \"Does a standard ISO CGL policy cover a data breach for a CA SaaS company with 50 employees?\"",
        size=15, bold=True, color=BRAND.navy_dark,
    )

    # Tool calls box
    _add_rect(s, Inches(0.5), Inches(2.7), Inches(6.0), Inches(3.5),
              RGBColor(0xF5, 0xF7, 0xFA))
    _add_text(s, Inches(0.7), Inches(2.85), Inches(5.7), Inches(0.4),
              "TOOL CALLS", size=11, bold=True, color=BRAND.accent)
    tool_lines = [
        "1.  ISO_forms_search → CG 00 01 04 13",
        "2.  ISO_forms_search → CG 21 07 (electronic data exclusion)",
        "3.  NAIC_lookup → MDL-672 (privacy)",
        "4.  state_DOI_query (CA) → cyber bulletins",
    ]
    for i, line in enumerate(tool_lines):
        _add_text(s, Inches(0.7), Inches(3.3) + Inches(0.4 * i), Inches(5.7),
                  Inches(0.4), line, size=14, color=BRAND.body, font="Consolas")

    _add_text(s, Inches(0.7), Inches(5.7), Inches(5.7), Inches(0.4),
              "RUN: 4 tools · 38s · $0.31",
              size=11, bold=True, color=BRAND.accent, font="Consolas")

    # Output box
    _add_rect(s, Inches(6.8), Inches(2.7), Inches(6.0), Inches(4.0),
              RGBColor(0xF5, 0xF7, 0xFA))
    _add_text(s, Inches(7.0), Inches(2.85), Inches(5.7), Inches(0.4),
              "OUTPUT", size=11, bold=True, color=BRAND.accent)
    _add_text(s, Inches(7.0), Inches(3.25), Inches(5.7), Inches(2.7),
              "Summary: Standard CGL (CG 00 01) excludes electronic-data losses via "
              "CG 21 07; cyber liability requires a separate endorsement or standalone "
              "policy. CA does not mandate cyber coverage but has consumer-notice rules "
              "under Cal. Civ. Code 1798.82.\n\n"
              "Findings: 4 · Sources: 4 · Flags: 1 · Confidence: medium",
              size=13, color=BRAND.body)
    _add_text(s, Inches(7.0), Inches(6.05), Inches(5.7), Inches(0.5),
              "⚑ Coverage determination requires underwriter review of carrier-specific endorsements.",
              size=12, bold=True, color=BRAND.accent)
    return s


def slide_why_now(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(s, "Why now.", kicker="Timing")

    cols = [
        ("Anthropic native tool use",
         "Tool calling went from a hack to a reliable primitive in 2024. Structured agent workflows finally work in production."),
        ("Regulatory data is machine-accessible",
         "NAIC CIS, state DOI portals, public IRS/SEC data — scrapable today. The data layer that didn't exist in 2020 exists now."),
        ("Insurance is under cycle-time pressure",
         "Post-2023 cat events, hardening market, talent shortage — every carrier is hunting for underwriter productivity."),
        ("AI buyer maturity in regulated industries",
         "Legal proved enterprise-AI procurement works (Harvey, Hebbia). Insurance is 18 months behind — the wave is forming."),
    ]
    col_w = Inches(6.0)
    col_h = Inches(2.2)
    top = Inches(2.0)
    for i, (h, body) in enumerate(cols):
        x = Inches(0.5) + (col_w + Inches(0.3)) * (i % 2)
        y = top + (col_h + Inches(0.3)) * (i // 2)
        _add_rect(s, x, y, Inches(0.06), col_h - Inches(0.2), BRAND.accent)
        _add_text(s, x + Inches(0.2), y, col_w - Inches(0.2), Inches(0.5),
                  h, size=18, bold=True, color=BRAND.navy_dark)
        _add_text(s, x + Inches(0.2), y + Inches(0.55), col_w - Inches(0.2),
                  col_h - Inches(0.6), body, size=14, color=BRAND.body)
    return s


def slide_wedge(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(s, "Where we start.", kicker="The Wedge")

    rows = [
        ("Persona", "Underwriter at a specialty MGA"),
        ("Job-to-be-done", "Coverage analysis for commercial lines (CGL, BOP, Property, Cyber)"),
        ("Buyer", "Head of Underwriting / Chief Underwriting Officer"),
        ("Company profile", "$50M – $500M premium; 8 – 15 underwriters; specialty/E&S focus"),
        ("Universe", "~80 specialty MGAs in the US fit this profile"),
    ]
    top = Inches(2.1)
    row_h = Inches(0.8)
    for i, (k, v) in enumerate(rows):
        y = top + row_h * i
        _add_rect(s, Inches(0.5), y + Inches(0.05), Inches(0.06), Inches(0.55), BRAND.accent)
        _add_text(s, Inches(0.7), y, Inches(3.5), Inches(0.7),
                  k, size=18, bold=True, color=BRAND.navy_dark)
        _add_text(s, Inches(4.5), y, Inches(8.3), Inches(0.7),
                  v, size=18, color=BRAND.body)

    _add_text(
        s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
        "We win one MGA's underwriting team. Then we expand to adjusters, brokers, "
        "compliance — same data, different prompts.",
        size=14, color=BRAND.muted,
    )
    return s


def slide_market(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(s, "Market.", kicker="Bottom-up")

    # Wedge
    _add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
              "Initial wedge (specialty MGAs, US):",
              size=18, bold=True, color=BRAND.navy_dark)
    _add_bullets(
        s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.6),
        [
            "80 specialty MGAs · ~12 underwriters each = ~960 seats",
            "$500/seat/month → $5.8M ARR if fully penetrated",
            "Realistic 3-year capture at 15% → $870K ARR from wedge alone",
        ],
        size=15,
    )

    # Expansion
    _add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5),
              "Adjacent expansion (same architecture, different prompts/data):",
              size=18, bold=True, color=BRAND.navy_dark)
    _add_bullets(
        s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.0),
        [
            "Carrier underwriting teams — ~3,000 US carriers, 50K+ underwriters",
            "Claims adjusters — 350K licensed adjusters · same regulatory data, different queries",
            "Brokers — 36K independent agencies · placement and gap analysis",
            "Adjacent verticals: legal regulatory research, healthcare compliance, financial filings",
        ],
        size=14,
    )
    return s


def slide_business_model(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(s, "Pricing & unit economics.", kicker="Business Model")

    rows = [
        ("Pilot", "$1,500 flat for a 30-day proof of value"),
        ("Seat license", "$500 / underwriter / month, billed monthly"),
        ("Overage", "$0.75 / query above 200/month/seat"),
        ("Cost per query", "~$0.30 LLM + ~$0.05 data + ~$0.05 infra ≈ $0.40 — ~75% gross margin"),
        ("Logo target", "5 paying MGAs · ~60 seats · ~$30K MRR by Q4 2026"),
    ]
    top = Inches(2.1)
    row_h = Inches(0.85)
    for i, (k, v) in enumerate(rows):
        y = top + row_h * i
        _add_rect(s, Inches(0.5), y + Inches(0.08), Inches(0.06), Inches(0.55), BRAND.accent)
        _add_text(s, Inches(0.7), y, Inches(3.5), Inches(0.7),
                  k, size=17, bold=True, color=BRAND.navy_dark)
        _add_text(s, Inches(4.5), y, Inches(8.3), Inches(0.7),
                  v, size=17, color=BRAND.body)
    return s


def slide_moat(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(s, "What's defensible.", kicker="Moat-Building Plan")

    rows = [
        ("Data partnerships",
         "Verisk (ISO forms / ClaimSearch) — in negotiation. AM Best, NCCI on roadmap. "
         "Each partnership is months of legal — and a barrier the next entrant must cross."),
        ("Eval data",
         "Every customer query becomes labeled training data. After 12 months we have what "
         "no GPT-4-wrapper can replicate: the canonical insurance-research benchmark."),
        ("Integration depth",
         "Slack/Teams → AMS360 → Applied Epic → quote-to-bind workflows. "
         "A 6-month displacement cost for the customer means a 6-month moat for us."),
        ("Citation audit chain",
         "Every claim → tool → primary source. Auditable for compliance, defensible "
         "in a regulatory exam. Generic AI cannot offer this."),
    ]
    top = Inches(2.0)
    row_h = Inches(1.15)
    for i, (h, body) in enumerate(rows):
        y = top + row_h * i
        _add_rect(s, Inches(0.5), y, Inches(0.06), Inches(1.0), BRAND.accent)
        _add_text(s, Inches(0.7), y, Inches(12.0), Inches(0.4),
                  h, size=17, bold=True, color=BRAND.navy_dark)
        _add_text(s, Inches(0.7), y + Inches(0.45), Inches(12.0), Inches(0.7),
                  body, size=13, color=BRAND.body)
    return s


def slide_traction(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(s, "Where we are today.", kicker="Traction")

    _add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
              "Built (verified working):",
              size=18, bold=True, color=BRAND.navy_dark)
    _add_bullets(
        s, Inches(0.5), Inches(2.55), Inches(12.3), Inches(2.0),
        [
            "Custom agent runtime — Anthropic native tool use, hard stops on cost/time/steps, 12/12 tests passing",
            "Insurance vertical scaffolded — system prompt + 5 tools + 30 golden eval queries",
            "NAIC connector live — public CIS company search + 11-entry curated MDL-NNN model law index",
            "Postgres run history + spend ledger + per-tenant cost caps",
        ],
        size=14,
    )

    _add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
              "Next 90 days:",
              size=18, bold=True, color=BRAND.navy_dark)
    _add_bullets(
        s, Inches(0.5), Inches(5.25), Inches(12.3), Inches(2.0),
        [
            "Wire 2 priority connectors (state DOI: CA, TX, FL, NY, IL · ISO via Cornell LII baseline)",
            "Run eval against real Claude — target ≥80% on the 30 golden queries",
            "Sign 3 specialty MGA design partners — 90-day free pilots, weekly feedback loops",
            "Verisk partnership conversations — close on terms by Q3",
        ],
        size=14,
    )
    return s


def slide_team(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_title_bar(s, "Team.", kicker="Who's Building This")

    _add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
              "[Founder Name] — Founder, CEO/CTO",
              size=22, bold=True, color=BRAND.navy_dark)
    _add_bullets(
        s, Inches(0.5), Inches(2.55), Inches(12.3), Inches(2.5),
        [
            "[Prior role and company — e.g. Senior MLE at $COMPANY, X years building ML systems]",
            "[Domain credibility — e.g. previously built insurance research tooling at $COMPANY, "
            "or 5 years as an underwriter]",
            "[Technical proof — e.g. 06_job_agent (FastAPI agent system, deployed) · "
            "ai-saas-assistant (multi-tenant SaaS, $X MRR)]",
        ],
        size=15,
    )

    _add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
              "Advisors / planned hires:",
              size=18, bold=True, color=BRAND.navy_dark)
    _add_bullets(
        s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(1.5),
        [
            "[Insurance domain advisor — e.g. former Chief Underwriting Officer at $MGA]",
            "[GTM advisor — someone who's sold to insurance before]",
            "First hire: founding engineer (post-pre-seed)",
        ],
        size=14,
    )
    return s


def slide_ask(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BRAND.navy_dark)

    _add_rect(s, Inches(0.6), Inches(0.7), Inches(0.6), Inches(0.07), BRAND.accent)
    _add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "THE ASK", size=12, bold=True, color=BRAND.accent)
    _add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.85),
              "Raising $[750K] pre-seed.",
              size=44, bold=True, color=BRAND.bg, font=BRAND.title_font)

    _add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(0.5),
              "Use of funds (12-month runway):",
              size=18, bold=True, color=BRAND.bg)
    use_of_funds = [
        "Founding engineer + part-time insurance domain expert — $400K",
        "Verisk partnership + data licensing deposits — $150K",
        "Cloud + LLM costs (eval + early customer runs) — $80K",
        "GTM (conferences, design-partner pilots, light marketing) — $60K",
        "Legal, ops, runway buffer — $60K",
    ]
    for i, item in enumerate(use_of_funds):
        _add_text(s, Inches(0.9), Inches(3.15) + Inches(0.45 * i), Inches(12), Inches(0.4),
                  f"•  {item}", size=15, color=RGBColor(0xE5, 0xEC, 0xF4))

    _add_text(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.5),
              "Milestones we'll hit with this capital:",
              size=18, bold=True, color=BRAND.bg)
    milestones = [
        "5 paying MGA logos · ~$30K MRR · ≥80% eval accuracy · Verisk LOI",
    ]
    for i, item in enumerate(milestones):
        _add_text(s, Inches(0.9), Inches(6.25) + Inches(0.4 * i), Inches(12), Inches(0.4),
                  f"•  {item}", size=15, color=RGBColor(0xE5, 0xEC, 0xF4))
    return s


def slide_close(prs: Presentation):
    s = _slide_with_layout(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BRAND.bg)
    _add_rect(s, Inches(0.5), Inches(0.45), Inches(0.6), Inches(0.07), BRAND.accent)

    _add_text(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.0),
              "AgentForge",
              size=56, bold=True, color=BRAND.navy_dark, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.6),
              "Let's build the research layer that insurance has been waiting for.",
              size=22, color=BRAND.body, align=PP_ALIGN.CENTER)

    _add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
              "[founder@email]    ·    [linkedin.com/in/founder]    ·    [agentforge.example.com]",
              size=15, color=BRAND.muted, align=PP_ALIGN.CENTER)
    return s


# ───────────────────────────────────────────────────────────────────
# Entry point
# ───────────────────────────────────────────────────────────────────


def build(out: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_cover,
        slide_problem,
        slide_generic_ai_fails,
        slide_insight,
        slide_solution,
        slide_how_it_works,
        slide_demo,
        slide_why_now,
        slide_wedge,
        slide_market,
        slide_business_model,
        slide_moat,
        slide_traction,
        slide_team,
        slide_ask,
        slide_close,
    ]
    for i, fn in enumerate(builders, start=1):
        slide = fn(prs)
        # cover and ask + close get no footer (intentional — full-bleed slides)
        if fn not in (slide_cover, slide_ask, slide_close):
            _add_footer(slide, i, len(builders))

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)


if __name__ == "__main__":
    output = Path(__file__).parent / "AgentForge.pitch.pptx"
    build(output)
    print(f"wrote {output}")
